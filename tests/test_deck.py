"""El deck impone la gramática de consultoría, no solo dibuja láminas."""

import matplotlib

matplotlib.use("Agg")

import pandas as pd
import pytest
from pptx import Presentation
from pptx.util import Inches

from sumakit import plots
from sumakit.deck import Deck, TituloNoAccionable

HALLAZGO = "Dos tercios de la cartera realizó tres transacciones o menos en el mes"


@pytest.fixture
def figura(tmp_path):
    fig = plots.ranking(pd.Series({"lunes": 3.0, "martes": 5.0, "miércoles": 1.0}))
    ruta = tmp_path / "fig.png"
    fig.savefig(ruta)
    return ruta


# --- la regla que justifica el módulo ---------------------------------------

def test_rechaza_un_titulo_que_es_etiqueta():
    """"Resultados" obliga al lector a deducir; un hallazgo se lo entrega."""
    with pytest.raises(TituloNoAccionable, match="etiqueta"):
        Deck("d").finding("Resultados")


def test_acepta_un_titulo_de_accion():
    d = Deck("d")
    d.finding(HALLAZGO)
    assert len(d) == 1


def test_el_mensaje_explica_como_arreglarlo():
    with pytest.raises(TituloNoAccionable) as e:
        Deck("d").finding("Conclusiones")
    assert "frase completa" in str(e.value) and "strict=False" in str(e.value)


def test_strict_false_es_la_salida_de_emergencia():
    d = Deck("d", strict=False)
    d.finding("Resultados")
    assert len(d) == 1


def test_la_regla_tambien_aplica_a_las_tablas():
    with pytest.raises(TituloNoAccionable):
        Deck("d").table("Datos", pd.DataFrame({"a": [1]}))


# --- estructura -------------------------------------------------------------

def test_las_laminas_son_16_9():
    d = Deck("d")
    assert d.prs.slide_width == Inches(13.333)
    assert d.prs.slide_height == Inches(7.5)


def test_construye_el_arco_completo(figura):
    d = Deck("Título", subtitle="sub", footer="SUMADOTS")
    d.cover()
    d.agenda(["Negocio", "Datos", "Hallazgos", "Recomendaciones"])
    d.section("Principales hallazgos")
    d.finding(HALLAZGO, kicker="Principales hallazgos", image=figura)
    d.closing("Gracias")
    assert len(d) == 5


def test_la_figura_queda_embebida(figura, tmp_path):
    d = Deck("d")
    d.finding(HALLAZGO, image=figura)
    ruta = d.save(tmp_path / "x.pptx")
    import zipfile
    with zipfile.ZipFile(ruta) as z:
        assert [n for n in z.namelist() if n.startswith("ppt/media/")], \
            "la figura no viajó dentro del pptx"


def test_figura_inexistente_falla_claro():
    with pytest.raises(FileNotFoundError, match="figura"):
        Deck("d").finding(HALLAZGO, image="no/existe.png")


def test_los_callouts_agregan_formas(figura):
    d = Deck("d")
    s_sin = d.finding(HALLAZGO, image=figura)
    n_sin = len(s_sin.shapes)
    s_con = d.finding(HALLAZGO, image=figura,
                      callouts=[("una anotación", 0.5, 0.1), ("otra", 0.1, 0.6)])
    assert len(s_con.shapes) == n_sin + 2


def test_el_titulo_de_accion_aparece_en_la_lamina(figura):
    d = Deck("d")
    s = d.finding(HALLAZGO, kicker="Hallazgos", image=figura)
    textos = " ".join(f.text_frame.text for f in s.shapes if f.has_text_frame)
    assert HALLAZGO in textos
    assert "Hallazgos" in textos


def test_la_tabla_lleva_encabezado_y_filas():
    datos = pd.DataFrame({"a": [1, 2], "b": [3.5, 4.5]}, index=["x", "y"])
    d = Deck("d")
    s = d.table(HALLAZGO, datos)
    tabla = next(f.table for f in s.shapes if f.has_table)
    assert tabla.cell(0, 1).text == "a"
    assert tabla.cell(1, 0).text == "x"


def test_el_archivo_guardado_se_puede_reabrir(tmp_path):
    d = Deck("d")
    d.cover()
    ruta = d.save(tmp_path / "sub/carpeta/x.pptx")
    assert ruta.exists()
    assert len(Presentation(str(ruta)).slides) == 1


def test_usa_el_color_de_acento_del_tema():
    from sumakit import theme
    d = Deck("d", palette=theme.LIGHT)
    assert str(d._acento) == theme.LIGHT.categorical[0].lstrip("#").upper()


# --- regresiones de lo que salió feo en el primer deck ----------------------

def _fuentes(slide):
    return {r.font.name
            for sh in slide.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs}


def test_todo_el_texto_lleva_fuente_explicita(figura):
    """Sin font.name, PowerPoint cae en Calibri y el tema no se aplica."""
    d = Deck("Un título de portada suficientemente largo", subtitle="s", footer="F")
    laminas = [d.cover(), d.agenda(["uno", "dos"]), d.section("Sección"),
               d.finding(HALLAZGO, kicker="k", image=figura,
                         callouts=[("nota de ejemplo", 0.5, 0.1)])]
    for s in laminas:
        assert None not in _fuentes(s), "hay texto sin tipografía fijada"


def test_la_tabla_no_usa_el_estilo_bandeado_de_office():
    d = Deck("d")
    s = d.table(HALLAZGO, pd.DataFrame({"a": [1, 2]}, index=["x", "y"]))
    tabla = next(f.table for f in s.shapes if f.has_table)
    assert tabla.first_row is False and tabla.horz_banding is False


def test_el_encabezado_de_la_tabla_usa_el_acento():
    d = Deck("d")
    s = d.table(HALLAZGO, pd.DataFrame({"a": [1]}, index=["x"]))
    tabla = next(f.table for f in s.shapes if f.has_table)
    assert tabla.cell(0, 0).fill.fore_color.rgb == d._acento


def test_la_figura_ocupa_el_ancho_disponible(figura):
    """Escalar solo por altura dejaba media lámina vacía."""
    d = Deck("d")
    s = d.finding(HALLAZGO, image=figura)
    img = next(sh for sh in s.shapes if sh.shape_type == 13)
    ancho = img.width / 914400
    alto = img.height / 914400
    assert ancho > 8.0, f"la figura quedó angosta: {ancho:.2f} pulgadas"
    assert alto <= 5.11, "la figura se sale del área de contenido"


def test_la_figura_queda_centrada(figura):
    d = Deck("d")
    s = d.finding(HALLAZGO, image=figura)
    img = next(sh for sh in s.shapes if sh.shape_type == 13)
    izq = img.left / 914400
    der = 13.333 - (img.left + img.width) / 914400
    assert abs(izq - der) < 0.02, "los márgenes laterales no coinciden"


def test_la_figura_conserva_su_proporcion(figura):
    from PIL import Image
    with Image.open(figura) as im:
        proporcion = im.size[0] / im.size[1]
    d = Deck("d")
    img = next(sh for sh in d.finding(HALLAZGO, image=figura).shapes
               if sh.shape_type == 13)
    assert abs(img.width / img.height - proporcion) < 0.01


# --- lo que se tomó del deck de referencia ---------------------------------

def test_el_titulo_resalta_lo_marcado_con_asteriscos():
    """La segunda mitad en otro color dirige la lectura sin agregar elementos."""
    d = Deck("d")
    s = d.finding("Más del 40% son cero, lo que **bloquea el log-ratio**")
    caja = next(sh for sh in s.shapes
                if sh.has_text_frame and "cero" in sh.text_frame.text)
    runs = caja.text_frame.paragraphs[0].runs
    assert len(runs) == 2
    assert runs[0].font.color.rgb == d._tinta
    assert runs[1].font.color.rgb == d._acento and runs[1].font.bold


def test_los_asteriscos_no_cuentan_como_palabras():
    with pytest.raises(TituloNoAccionable):
        Deck("d").finding("**Resultados**")


def test_la_lamina_hero_pone_el_numero_en_grande():
    d = Deck("d")
    s = d.hero("16.705", "clientes con una sola transacción", kicker="Hallazgos")
    tamaños = {r.font.size.pt for sh in s.shapes if sh.has_text_frame
               for p in sh.text_frame.paragraphs for r in p.runs}
    assert max(tamaños) >= 60, "el número héroe debe dominar la lámina"


def test_el_hero_usa_el_acento_para_la_cifra():
    d = Deck("d")
    s = d.hero("42%", "de la cartera")
    caja = next(sh for sh in s.shapes if sh.has_text_frame and "42%" in sh.text_frame.text)
    assert caja.text_frame.paragraphs[0].runs[0].font.color.rgb == d._acento


def test_el_tema_oscuro_cambia_el_fondo():
    from sumakit import theme
    claro, oscuro = Deck("d", palette=theme.LIGHT), Deck("d", palette=theme.DARK)
    assert claro._fondo != oscuro._fondo
    assert oscuro.es_oscuro and not claro.es_oscuro
