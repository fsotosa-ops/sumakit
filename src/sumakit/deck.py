"""Decks de consultoría generados por código.

Quarto produce presentaciones correctas pero con su propio layout: no hace el
formato de consultoría —etiqueta de sección, título de acción, una figura
anotada con cajas punteadas—. Este módulo sí, sobre `python-pptx`, y el
resultado es un `.pptx` de verdad: se abre en PowerPoint o Google Slides y se
retoca ahí. Esa es la interfaz; no hay que construir ninguna.

Lo que la API impone, y es su razón de ser: **no se puede crear una lámina de
hallazgo sin un título de acción**. Un título como "Resultados" es una
etiqueta; "Dos tercios de la cartera realizó tres transacciones o menos" es un
hallazgo. La primera obliga al lector a deducir; la segunda le entrega la
conclusión. Es la regla de la Pirámide de Minto, y aquí es un `ValueError`.

    from sumakit.deck import Deck

    d = Deck("Segmentación de tarjetahabientes", subtitle="Banco · agosto 2026")
    d.cover()
    d.agenda(["Negocio", "Datos", "Hallazgos", "Recomendaciones"])
    d.section("Principales hallazgos")
    d.finding(
        "Dos tercios de la cartera realizó tres transacciones o menos en el mes",
        image="figuras/concentracion.png",
        callouts=[("Con una transacción no hay patrón, hay un evento", 0.58, 0.12)],
    )
    d.save("entrega/negocio.pptx")
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import pandas as pd

from . import theme


def pd_api_numerica(serie) -> bool:
    return pd.api.types.is_numeric_dtype(serie)

__all__ = ["Deck", "TituloNoAccionable"]

# Geometría en pulgadas, 16:9. Sale del mockup validado a 96 px/pulgada.
_ANCHO, _ALTO = 13.333, 7.5
_MARGEN = 0.5
_ANCHO_UTIL = _ANCHO - 2 * _MARGEN

_KICKER_Y, _KICKER_H = 0.35, 0.28
_TITULO_Y, _TITULO_H = 0.63, 1.05
_CUERPO_Y = 1.77
_CUERPO_H = 5.10
_PIE_H = 0.30

_MIN_PALABRAS_TITULO = 5
# Con anotaciones al lado, la figura toma esta fracción del ancho útil.
_FRACCION_FIGURA = 0.66

# Sin esto PowerPoint usa Calibri y el tema no se aplica a nada.
_FUENTE = "Helvetica Neue"


class TituloNoAccionable(ValueError):
    """El título describe una categoría en vez de afirmar un hallazgo."""


def _rgb(hexa: str) -> RGBColor:
    return RGBColor.from_string(hexa.lstrip("#").upper())


class Deck:
    """Constructor de presentaciones con la gramática de consultoría."""

    def __init__(
        self,
        title: str,
        *,
        subtitle: str = "",
        footer: str = "",
        palette: theme.Palette | None = None,
        strict: bool = True,
    ) -> None:
        self.title = title
        self.subtitle = subtitle
        self.footer = footer
        self.palette = palette or theme.active()
        self.strict = strict

        self.prs = Presentation()
        self.prs.slide_width = Inches(_ANCHO)
        self.prs.slide_height = Inches(_ALTO)
        self._acento = _rgb(self.palette.categorical[0])
        self._tinta = _rgb(self.palette.text_primary)
        self._suave = _rgb(self.palette.text_secondary)
        self._fondo = _rgb(self.palette.surface)

    # --- primitivas ---------------------------------------------------------

    def _lamina(self):
        """Lámina en blanco: el layout lo ponemos nosotros, no la plantilla."""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fondo = s.background.fill
        fondo.solid()
        fondo.fore_color.rgb = self._fondo
        return s

    def _texto(self, s, x, y, w, h, texto, *, size=14, color=None, bold=False,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, espaciado=1.15):
        caja = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = caja.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = espaciado
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(size)
        run.font.name = _FUENTE
        run.font.bold = bold
        run.font.color.rgb = color or self._tinta
        return caja

    def _banda(self, s, y, h):
        from pptx.enum.shapes import MSO_SHAPE
        forma = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Inches(y),
                                   self.prs.slide_width, Inches(h))
        forma.fill.solid()
        forma.fill.fore_color.rgb = self._acento
        forma.line.fill.background()
        forma.shadow.inherit = False
        return forma

    def _pie(self, s):
        if self.footer:
            self._texto(s, _ANCHO - _MARGEN - 3.0, _ALTO - _PIE_H - 0.22, 3.0, _PIE_H,
                        self.footer, size=9, color=self._suave, align=PP_ALIGN.RIGHT)

    def _encabezado(self, s, kicker: str, titulo: str):
        if kicker:
            self._texto(s, _MARGEN, _KICKER_Y, _ANCHO_UTIL, _KICKER_H,
                        kicker, size=11, color=self._suave)
        self._texto(s, _MARGEN, _TITULO_Y, _ANCHO_UTIL, _TITULO_H,
                    titulo, size=18, espaciado=1.2)

    def _validar_titulo(self, titulo: str) -> None:
        if not self.strict:
            return
        if len(titulo.split()) < _MIN_PALABRAS_TITULO:
            raise TituloNoAccionable(
                f"«{titulo}» parece una etiqueta, no un hallazgo. Un título de "
                "acción afirma la conclusión en una frase completa: no "
                "«Resultados» sino «Dos tercios de la cartera realizó tres "
                "transacciones o menos». Si es a propósito, usa Deck(strict=False)."
            )

    # --- láminas ------------------------------------------------------------

    def cover(self):
        """Portada: bandas de acento arriba y abajo, título a la izquierda."""
        s = self._lamina()
        self._banda(s, 0, 0.15)
        self._banda(s, _ALTO - 0.55, 0.55)
        self._texto(s, _MARGEN, 3.1, _ANCHO_UTIL * 0.72, 1.4, self.title,
                    size=30, espaciado=1.2)
        if self.subtitle:
            self._texto(s, _MARGEN, 4.35, _ANCHO_UTIL * 0.72, 0.5, self.subtitle,
                        size=14, color=self._suave)
        return s

    def agenda(self, items: list[str], *, title: str = "Agenda"):
        """Índice: una banda con los puntos centrados, como el deck de referencia."""
        s = self._lamina()
        self._texto(s, _MARGEN, _KICKER_Y, _ANCHO_UTIL, _KICKER_H,
                    title, size=13, color=self._suave)
        alto = 0.52 * len(items) + 0.6
        self._banda(s, (_ALTO - alto) / 2, alto)
        caja = s.shapes.add_textbox(Inches(_MARGEN), Inches((_ALTO - alto) / 2 + 0.35),
                                    Inches(_ANCHO_UTIL), Inches(alto - 0.7))
        tf = caja.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.6
            run = p.add_run()
            run.text = item
            run.font.size = Pt(17)
            run.font.name = _FUENTE
            run.font.color.rgb = RGBColor.from_string("FFFFFF")
        self._pie(s)
        return s

    def section(self, name: str):
        """Separador de sección: banda ancha con el nombre."""
        s = self._lamina()
        self._banda(s, 2.9, 1.5)
        self._texto(s, _MARGEN, 3.25, _ANCHO_UTIL, 0.9, name,
                    size=26, color=RGBColor.from_string("FFFFFF"),
                    anchor=MSO_ANCHOR.MIDDLE)
        return s

    def finding(
        self,
        title: str,
        *,
        kicker: str = "",
        image: str | Path | None = None,
        body: str | None = None,
        callouts: list[tuple] | None = None,
    ):
        """Lámina de hallazgo: la unidad del deck.

        `title` debe ser un título de acción — la conclusión en una frase.

        `callouts` son cajas de anotación sobre el contenido, cada una
        `(texto, x, y)` con `x` e `y` como fracción del área de contenido, de
        modo que la posición no depende del tamaño de la lámina.
        """
        self._validar_titulo(title)
        s = self._lamina()
        self._encabezado(s, kicker, title)

        anotaciones = list(callouts or [])
        # Con anotaciones, la figura cede el tercio derecho: encima taparía
        # los datos, que es exactamente lo que un callout no debe hacer.
        fraccion = _FRACCION_FIGURA if (image is not None and anotaciones) else 1.0

        if image is not None:
            self._imagen(s, image, fraccion=fraccion)
        elif body:
            self._texto(s, _MARGEN, _CUERPO_Y, _ANCHO_UTIL, _CUERPO_H, body, size=14)

        self._anotar(s, anotaciones, fraccion)
        self._pie(s)
        return s

    def _anotar(self, s, anotaciones, fraccion):
        """Coloca las anotaciones apiladas en la columna libre de la derecha.

        Acepta texto suelto —y entonces las posiciona solas— o la forma
        `(texto, x, y)` para cuando quieras ubicarlas a mano.
        """
        if not anotaciones:
            return
        if fraccion >= 1.0:
            for a in anotaciones:
                texto, fx, fy = a if isinstance(a, (tuple, list)) else (a, 0.6, 0.05)
                self._callout(s, texto, fx, fy, ancho=_ANCHO_UTIL * 0.3)
            return

        ancho = _ANCHO_UTIL * (1 - fraccion) - 0.25
        x = _MARGEN + _ANCHO_UTIL * fraccion + 0.25
        y = _CUERPO_Y + 0.15
        for a in anotaciones:
            texto = a[0] if isinstance(a, (tuple, list)) else a
            forma = self._callout_abs(s, texto, x, y, ancho)
            y += forma.height / 914400 + 0.28

    def _imagen(self, s, image, *, fraccion: float = 1.0):
        """Encaja la figura en el área de contenido conservando su proporción.

        Escalar solo por altura dejaba la figura a la izquierda y media lámina
        vacía. Aquí se toma el factor que más limita y se centra el sobrante.
        """
        from PIL import Image as _Image

        ruta = Path(image)
        if not ruta.exists():
            raise FileNotFoundError(f"no encuentro la figura: {ruta}")
        with _Image.open(ruta) as im:
            ancho_px, alto_px = im.size
        disponible = _ANCHO_UTIL * fraccion
        proporcion = ancho_px / alto_px
        ancho, alto = disponible, disponible / proporcion
        if alto > _CUERPO_H:
            alto, ancho = _CUERPO_H, _CUERPO_H * proporcion
        x = _MARGEN + (disponible - ancho) / 2
        return s.shapes.add_picture(str(ruta), Inches(x), Inches(_CUERPO_Y),
                                    width=Inches(ancho), height=Inches(alto))

    def _callout(self, s, texto: str, fx: float, fy: float, ancho: float = 2.9):
        return self._callout_abs(s, texto, _MARGEN + fx * _ANCHO_UTIL,
                                 _CUERPO_Y + fy * _CUERPO_H, ancho)

    def _callout_abs(self, s, texto: str, x: float, y: float, ancho: float):
        from pptx.enum.shapes import MSO_SHAPE
        caracteres_por_linea = max(18, int(ancho * 15))
        lineas = len(texto) // caracteres_por_linea + 1
        alto = 0.30 + 0.21 * lineas
        forma = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                   Inches(ancho), Inches(alto))
        forma.fill.solid()
        forma.fill.fore_color.rgb = self._fondo
        forma.line.color.rgb = self._acento
        forma.line.width = Pt(1.25)
        forma.line.dash_style = 4  # punteado, como en el deck de referencia
        forma.shadow.inherit = False
        tf = forma.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.12)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(11)
        run.font.name = _FUENTE
        run.font.color.rgb = self._tinta
        return forma

    def table(self, title: str, data, *, kicker: str = "", max_rows: int = 12):
        """Lámina con una tabla, a partir de un DataFrame."""
        self._validar_titulo(title)
        s = self._lamina()
        self._encabezado(s, kicker, title)

        recorte = data.head(max_rows)
        filas, cols = recorte.shape[0] + 1, recorte.shape[1] + 1
        forma = s.shapes.add_table(filas, cols, Inches(_MARGEN), Inches(_CUERPO_Y),
                                   Inches(_ANCHO_UTIL), Inches(min(_CUERPO_H, 0.36 * filas)))
        tabla = forma.table
        # El estilo por defecto de Office pinta bandas azules. Se apaga y se
        # dibuja a mano: encabezado con el acento, filas sobre el fondo.
        tabla.first_row = False
        tabla.horz_banding = False
        tabla.cell(0, 0).text = str(recorte.index.name or "")
        for j, col in enumerate(recorte.columns, start=1):
            tabla.cell(0, j).text = str(col)
        for i, (idx, fila) in enumerate(recorte.iterrows(), start=1):
            tabla.cell(i, 0).text = str(idx)
            for j, valor in enumerate(fila, start=1):
                tabla.cell(i, j).text = f"{valor:,.4g}" if isinstance(valor, float) else str(valor)
        # Repartir el ancho por el largo del contenido: en partes iguales, el
        # nombre de una variable queda apretado y un número queda con aire.
        indices = [str(v) for v in recorte.index]
        pesos = [max(len(str(recorte.index.name or "")),
                     *(len(t) for t in indices)) if indices else 6]
        for col in recorte.columns:
            textos = [f"{v:,.4g}" if isinstance(v, float) else str(v) for v in recorte[col]]
            pesos.append(max(len(str(col)), *(len(t) for t in textos)) if textos
                         else len(str(col)))
        total = sum(pesos)
        for col_t, peso in zip(tabla.columns, pesos):
            col_t.width = Emu(int(Inches(_ANCHO_UTIL) * peso / total))

        numericas = {j for j, col in enumerate(recorte.columns, start=1)
                     if pd_api_numerica(recorte[col])}
        blanco = RGBColor.from_string("FFFFFF")
        for i, fila_t in enumerate(tabla.rows):
            for j, celda in enumerate(fila_t.cells):
                celda.fill.solid()
                celda.fill.fore_color.rgb = self._acento if i == 0 else self._fondo
                celda.margin_left = celda.margin_right = Inches(0.08)
                for p in celda.text_frame.paragraphs:
                    if j in numericas:
                        p.alignment = PP_ALIGN.RIGHT
                    for run in p.runs:
                        run.font.size = Pt(11)
                        run.font.name = _FUENTE
                        run.font.color.rgb = blanco if i == 0 else self._tinta
        self._pie(s)
        return s

    def closing(self, text: str = ""):
        s = self._lamina()
        self._banda(s, 0, 0.15)
        self._banda(s, _ALTO - 0.55, 0.55)
        if text:
            self._texto(s, _MARGEN, 3.3, _ANCHO_UTIL, 1.0, text, size=22)
        return s

    # --- salida -------------------------------------------------------------

    def save(self, path: str | Path) -> Path:
        ruta = Path(path)
        ruta.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(ruta))
        return ruta

    def __len__(self) -> int:
        return len(self.prs.slides)

    def __repr__(self) -> str:
        return f"Deck({self.title!r}, {len(self)} láminas)"
